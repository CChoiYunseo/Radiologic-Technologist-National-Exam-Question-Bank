#!/usr/bin/env python3
"""
Extract subject reference materials into page-preserving RAG chunks.

This script does not rewrite source files into secondary PDFs. It reads original
PDF/PPTX/DOCX/TXT files, preserves page or slide metadata, and writes JSON/JSONL
artifacts for later review, embedding, and DB import.
"""

from __future__ import annotations

import argparse
import hashlib
import io
import json
import re
from dataclasses import dataclass
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import fitz
import pdfplumber
from PIL import Image

try:
    import cv2
    import numpy as np
except Exception:  # pragma: no cover - optional runtime fallback
    cv2 = None
    np = None

try:
    import pytesseract
except Exception:  # pragma: no cover - checked at runtime
    pytesseract = None

try:
    from docx import Document
except Exception:  # pragma: no cover - optional runtime fallback
    Document = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover - optional runtime fallback
    Presentation = None


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_INPUT = PROJECT_ROOT / "materials" / "04_subject_references"
DEFAULT_OUTPUT = PROJECT_ROOT / "resources" / "extracted" / "subject_references"
SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt", ".md"}


@dataclass
class PageRecord:
    source_file: str
    source_path: str
    source_type: str
    page_or_slide: int
    content: str
    extraction_method: str
    extraction_quality: str
    ocr_confidence: float | None = None
    table_count: int = 0
    image_count: int = 0
    visual_image_count: int = 0
    drawing_count: int = 0


def now_iso() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat()


def stable_id(*parts: Any, length: int = 32) -> str:
    raw = "::".join(str(part) for part in parts)
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()[:length]


def clean_text(value: str) -> str:
    text = str(value or "").replace("\x00", " ")
    text = re.sub(r"[ \t\r\f\v]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines = [line.strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def content_quality(text: str, ocr_confidence: float | None = None) -> str:
    length = len(clean_text(text))
    if ocr_confidence is not None and ocr_confidence < 45:
        return "low"
    if length >= 500 and (ocr_confidence is None or ocr_confidence >= 70):
        return "high"
    if length >= 120:
        return "medium"
    return "low"


def table_to_text(table: list[list[Any]]) -> str:
    rows = []
    for row in table or []:
        cells = [clean_text(cell or "") for cell in row]
        if any(cells):
            rows.append(" | ".join(cells))
    return "\n".join(rows)


def extract_pdf_tables(pdf_path: Path, max_pages: int | None = None) -> dict[int, list[str]]:
    tables_by_page: dict[int, list[str]] = {}
    try:
        with pdfplumber.open(str(pdf_path)) as pdf:
            for index, page in enumerate(pdf.pages, start=1):
                if max_pages and index > max_pages:
                    break
                page_tables = []
                for table in page.extract_tables() or []:
                    rendered = table_to_text(table)
                    if rendered:
                        page_tables.append(rendered)
                if page_tables:
                    tables_by_page[index] = page_tables
    except Exception as exc:
        tables_by_page[0] = [f"[table extraction failed: {type(exc).__name__}: {exc}]"]
    return tables_by_page


def preprocess_for_ocr(image: Image.Image) -> Image.Image:
    if cv2 is None or np is None:
        return image
    array = np.array(image.convert("RGB"))
    gray = cv2.cvtColor(array, cv2.COLOR_RGB2GRAY)
    gray = cv2.medianBlur(gray, 3)
    thresholded = cv2.adaptiveThreshold(
        gray,
        255,
        cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
        cv2.THRESH_BINARY,
        35,
        11,
    )
    return Image.fromarray(thresholded)


def ocr_pdf_page(page: fitz.Page, dpi: int, lang: str) -> tuple[str, float | None]:
    if pytesseract is None:
        return "", None
    matrix = fitz.Matrix(dpi / 72, dpi / 72)
    pixmap = page.get_pixmap(matrix=matrix, alpha=False)
    image = Image.open(io.BytesIO(pixmap.tobytes("png")))
    image = preprocess_for_ocr(image)
    text = pytesseract.image_to_string(image, lang=lang, config="--psm 6")

    confidence: float | None = None
    try:
        data = pytesseract.image_to_data(image, lang=lang, config="--psm 6", output_type=pytesseract.Output.DICT)
        values = []
        for raw in data.get("conf", []):
            try:
                value = float(raw)
            except Exception:
                continue
            if value >= 0:
                values.append(value)
        if values:
            confidence = sum(values) / len(values)
    except Exception:
        confidence = None
    return clean_text(text), confidence


def count_pdf_images(page: fitz.Page, min_area_ratio: float, max_area_ratio: float) -> tuple[int, int]:
    page_area = page.rect.width * page.rect.height
    infos = page.get_image_info(xrefs=True)
    visual_count = 0
    for info in infos:
        bbox = fitz.Rect(info.get("bbox", (0, 0, 0, 0)))
        ratio = (bbox.width * bbox.height) / page_area if page_area else 0
        if min_area_ratio <= ratio <= max_area_ratio:
            visual_count += 1
    return len(infos), visual_count


def extract_pdf(path: Path, args: argparse.Namespace) -> list[PageRecord]:
    records: list[PageRecord] = []
    doc = fitz.open(str(path))
    max_pages = min(args.max_pages or len(doc), len(doc))
    tables_by_page = extract_pdf_tables(path, max_pages=max_pages) if args.extract_tables else {}

    for page_index in range(max_pages):
        page_number = page_index + 1
        page = doc[page_index]
        text = clean_text(page.get_text("text", sort=True))
        method = "pymupdf_text"
        confidence = None
        image_count = 0
        visual_image_count = 0
        drawing_count = 0
        if args.detect_visuals:
            image_count, visual_image_count = count_pdf_images(
                page,
                args.visual_image_min_area_ratio,
                args.visual_image_max_area_ratio,
            )
            drawing_count = len(page.get_drawings())

        if args.ocr and len(text) < args.ocr_min_chars:
            ocr_text, confidence = ocr_pdf_page(page, args.ocr_dpi, args.ocr_lang)
            if len(ocr_text) > len(text):
                text = ocr_text
                method = "tesseract_ocr"
            elif ocr_text:
                text = clean_text(text + "\n\n[OCR 보조 추출]\n" + ocr_text)
                method = "pymupdf_text+tesseract_ocr"

        page_tables = tables_by_page.get(page_number, [])
        if page_tables:
            text = clean_text(text + "\n\n[표 추출]\n" + "\n\n".join(page_tables))

        records.append(PageRecord(
            source_file=path.name,
            source_path=str(path.relative_to(PROJECT_ROOT)),
            source_type="pdf",
            page_or_slide=page_number,
            content=text,
            extraction_method=method,
            extraction_quality=content_quality(text, confidence),
            ocr_confidence=confidence,
            table_count=len(page_tables),
            image_count=image_count,
            visual_image_count=visual_image_count,
            drawing_count=drawing_count,
        ))
    doc.close()
    return records


def extract_pptx(path: Path, args: argparse.Namespace) -> list[PageRecord]:
    if Presentation is None:
        raise RuntimeError("python-pptx is not installed")
    prs = Presentation(str(path))
    records = []
    max_slides = min(args.max_pages or len(prs.slides), len(prs.slides))
    for index, slide in enumerate(prs.slides[:max_slides], start=1):
        parts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                parts.append(shape.text)
            if getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    parts.append(" | ".join(clean_text(cell.text) for cell in row.cells))
        text = clean_text("\n".join(parts))
        records.append(PageRecord(
            source_file=path.name,
            source_path=str(path.relative_to(PROJECT_ROOT)),
            source_type="pptx",
            page_or_slide=index,
            content=text,
            extraction_method="python-pptx",
            extraction_quality=content_quality(text),
        ))
    return records


def extract_docx(path: Path, args: argparse.Namespace) -> list[PageRecord]:
    if Document is None:
        raise RuntimeError("python-docx is not installed")
    doc = Document(str(path))
    parts = [paragraph.text for paragraph in doc.paragraphs if clean_text(paragraph.text)]
    for table in doc.tables:
        for row in table.rows:
            parts.append(" | ".join(clean_text(cell.text) for cell in row.cells))
    text = clean_text("\n".join(parts))
    return [PageRecord(
        source_file=path.name,
        source_path=str(path.relative_to(PROJECT_ROOT)),
        source_type="docx",
        page_or_slide=1,
        content=text,
        extraction_method="python-docx",
        extraction_quality=content_quality(text),
    )]


def extract_text_file(path: Path) -> list[PageRecord]:
    text = clean_text(path.read_text(encoding="utf-8", errors="ignore"))
    return [PageRecord(
        source_file=path.name,
        source_path=str(path.relative_to(PROJECT_ROOT)),
        source_type=path.suffix.lower().lstrip("."),
        page_or_slide=1,
        content=text,
        extraction_method="plain_text",
        extraction_quality=content_quality(text),
    )]


def extract_file(path: Path, args: argparse.Namespace) -> list[PageRecord]:
    suffix = path.suffix.lower()
    if suffix == ".pdf":
        return extract_pdf(path, args)
    if suffix == ".pptx":
        return extract_pptx(path, args)
    if suffix == ".docx":
        return extract_docx(path, args)
    if suffix in {".txt", ".md"}:
        return extract_text_file(path)
    return []


def split_into_chunks(text: str, max_chars: int, overlap_chars: int) -> list[str]:
    text = clean_text(text)
    if not text:
        return []
    paragraphs = re.split(r"\n\s*\n", text)
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        paragraph = clean_text(paragraph)
        if not paragraph:
            continue
        if len(paragraph) > max_chars:
            if current:
                chunks.append(current)
                current = ""
            start = 0
            while start < len(paragraph):
                end = min(start + max_chars, len(paragraph))
                chunks.append(paragraph[start:end].strip())
                if end == len(paragraph):
                    break
                start = max(0, end - overlap_chars)
            continue
        candidate = clean_text(current + "\n\n" + paragraph) if current else paragraph
        if len(candidate) <= max_chars:
            current = candidate
        else:
            chunks.append(current)
            prefix = current[-overlap_chars:] if overlap_chars and current else ""
            current = clean_text(prefix + "\n\n" + paragraph) if prefix else paragraph
    if current:
        chunks.append(current)
    return chunks


def discover_files(input_dir: Path) -> list[Path]:
    files = []
    for path in sorted(input_dir.rglob("*")):
        if path.is_file() and path.suffix.lower() in SUPPORTED_EXTENSIONS:
            if path.name.startswith(".") or path.name.upper() == "README.MD":
                continue
            files.append(path)
    return files


def build_document_record(path: Path, pages: list[PageRecord], created_at: str) -> dict[str, Any]:
    document_id = stable_id(path.relative_to(PROJECT_ROOT), path.stat().st_size, path.stat().st_mtime_ns)
    low_pages = [record.page_or_slide for record in pages if record.extraction_quality == "low"]
    text_pages = [record for record in pages if record.content]
    return {
        "document_id": document_id,
        "source_file": path.name,
        "source_path": str(path.relative_to(PROJECT_ROOT)),
        "source_type": path.suffix.lower().lstrip("."),
        "page_count": len(pages),
        "text_pages": len(text_pages),
        "low_quality_pages": low_pages,
        "extraction_status": "completed" if pages else "failed",
        "copyright_status": "unknown",
        "created_at": created_at,
    }


def page_to_chunks(document: dict[str, Any], page: PageRecord, args: argparse.Namespace, created_at: str) -> list[dict[str, Any]]:
    chunks = []
    has_visual_elements = bool(args.detect_visuals and (page.visual_image_count > 0 or page.drawing_count >= args.visual_drawing_threshold))
    visual_review_priority = "not_applicable"
    if has_visual_elements:
        visual_review_priority = "high" if page.visual_image_count > 0 or page.drawing_count >= args.visual_high_drawing_threshold else "medium"
    for offset, content in enumerate(split_into_chunks(page.content, args.chunk_chars, args.overlap_chars)):
        chunk_id = stable_id(document["document_id"], page.page_or_slide, offset, content)
        chunks.append({
            "chunk_id": chunk_id,
            "document_id": document["document_id"],
            "chunk_index": offset,
            "source_file": page.source_file,
            "source_type": page.source_type,
            "page_or_slide": page.page_or_slide,
            "exam_period": "",
            "subject": "",
            "field": "",
            "area": "",
            "sub_area": "",
            "learning_objective": "",
            "keywords": [],
            "content": content,
            "content_hash": stable_id(content, length=64),
            "token_estimate": max(1, len(content) // 3),
            "extraction_method": page.extraction_method,
            "extraction_quality": page.extraction_quality,
            "ocr_confidence": page.ocr_confidence,
            "table_count": page.table_count,
            "table_extraction_method": "pdfplumber_structured_table" if page.table_count else "",
            "image_count": page.image_count,
            "visual_image_count": page.visual_image_count,
            "drawing_count": page.drawing_count,
            "has_visual_elements": has_visual_elements,
            "visual_element_count": page.image_count + page.drawing_count,
            "visual_description": "",
            "visual_description_method": "",
            "visual_description_status": "needs_vision_llm" if has_visual_elements else "not_applicable",
            "visual_review_priority": visual_review_priority,
            "needs_review": page.extraction_quality != "high",
            "created_at": created_at,
        })
    return chunks


def write_json(path: Path, data: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def write_jsonl(path: Path, rows: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as f:
        for row in rows:
            f.write(json.dumps(row, ensure_ascii=False) + "\n")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description="Extract subject reference files into RAG-ready chunks.")
    parser.add_argument("--input", type=Path, default=DEFAULT_INPUT)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--chunk-chars", type=int, default=1400)
    parser.add_argument("--overlap-chars", type=int, default=180)
    parser.add_argument("--ocr-min-chars", type=int, default=120)
    parser.add_argument("--ocr-dpi", type=int, default=300)
    parser.add_argument("--ocr-lang", default="kor+eng")
    parser.add_argument("--visual-image-min-area-ratio", type=float, default=0.15)
    parser.add_argument("--visual-image-max-area-ratio", type=float, default=0.9)
    parser.add_argument("--visual-drawing-threshold", type=int, default=500)
    parser.add_argument("--visual-high-drawing-threshold", type=int, default=500)
    parser.add_argument("--no-ocr", dest="ocr", action="store_false")
    parser.add_argument("--extract-tables", dest="extract_tables", action="store_true")
    parser.add_argument("--detect-visuals", dest="detect_visuals", action="store_true")
    parser.add_argument("--max-pages", type=int, default=0, help="Limit pages/slides per file for smoke tests.")
    parser.set_defaults(ocr=True, extract_tables=False, detect_visuals=False)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    input_dir = args.input.resolve()
    output_dir = args.output.resolve()
    created_at = now_iso()

    files = discover_files(input_dir)
    documents: list[dict[str, Any]] = []
    chunks: list[dict[str, Any]] = []
    errors: list[dict[str, str]] = []

    for path in files:
        try:
            pages = extract_file(path, args)
            document = build_document_record(path, pages, created_at)
            documents.append(document)
            for page in pages:
                chunks.extend(page_to_chunks(document, page, args, created_at))
        except Exception as exc:
            errors.append({
                "source_file": path.name,
                "source_path": str(path.relative_to(PROJECT_ROOT)),
                "error": f"{type(exc).__name__}: {exc}",
            })

    report = {
        "created_at": created_at,
        "input_dir": str(input_dir),
        "output_dir": str(output_dir),
        "file_count": len(files),
        "document_count": len(documents),
        "chunk_count": len(chunks),
        "low_quality_chunk_count": sum(1 for chunk in chunks if chunk["extraction_quality"] == "low"),
        "needs_review_chunk_count": sum(1 for chunk in chunks if chunk["needs_review"]),
        "errors": errors,
        "notes": [
            "원문 파일은 수정하지 않는다.",
            "원문을 대체하는 2차 PDF를 생성하지 않는다.",
            "MVP 기본 실행에서는 텍스트 레이어 추출과 OCR 기반 텍스트 추출만 수행한다.",
            "표 구조화 추출은 --extract-tables 옵션을 켤 때만 수행한다.",
            "그림, 도식, 장비 구조도 후보 탐지는 --detect-visuals 옵션을 켤 때만 수행한다.",
            "출제범위와 학습목표는 자동 확정하지 않고 후속 매핑 단계에서 연결한다.",
        ],
    }

    write_json(output_dir / "documents.json", documents)
    write_jsonl(output_dir / "chunks.jsonl", chunks)
    write_json(output_dir / "extraction_report.json", report)
    print(json.dumps(report, ensure_ascii=False))


if __name__ == "__main__":
    main()
