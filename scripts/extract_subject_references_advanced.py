#!/usr/bin/env python3
"""Extract multimodal subject references into reviewable RAG chunks.

The script keeps the original teaching materials intact. It decomposes PDFs into
page-level text, structured table chunks, formula candidates, and visual
candidates. Heavy crop images are written outside /mnt/data by default, while
JSON/JSONL metadata stays in the project resources folder.
"""

from __future__ import annotations

import argparse
import hashlib
import io
import json
import re
import unicodedata
from collections import Counter, defaultdict
from datetime import datetime, timezone
from pathlib import Path
from typing import Any

import fitz
import pdfplumber
from PIL import Image

try:
    import cv2
    import numpy as np
except Exception:  # pragma: no cover
    cv2 = None
    np = None

try:
    import pytesseract
except Exception:  # pragma: no cover
    pytesseract = None

try:
    from docx import Document
except Exception:  # pragma: no cover
    Document = None

try:
    from pptx import Presentation
except Exception:  # pragma: no cover
    Presentation = None


PROJECT_ROOT = Path(__file__).resolve().parents[1]
DEFAULT_INPUT = PROJECT_ROOT / "materials" / "04_subject_references"
DEFAULT_OUTPUT = PROJECT_ROOT / "resources" / "extracted" / "subject_references_advanced"
DEFAULT_ASSET_OUTPUT = Path("/opt/app/extracted_assets/subject_references_advanced")
EXTERNAL_MATERIALS_ROOT = Path("/opt/app/materials")
SUPPORTED_EXTENSIONS = {".pdf", ".pptx", ".docx", ".txt", ".md"}

FORMULA_PATTERNS = [
    re.compile(r"[A-Za-z가-힣0-9)\]]\s*[=＝]\s*[-+*/×÷\w가-힣([{\ue000-\uf8ff]"),
    re.compile(r"[∑√πΩμθλαβγ∆Δ±≤≥≠∞∝]"),
    re.compile(r"\b(?:sin|cos|tan|log|ln)\b", re.IGNORECASE),
]


def now_iso() -> str:
    return datetime.now(timezone.utc).replace(microsecond=0).isoformat()


def nfc(value: Any) -> str:
    return unicodedata.normalize("NFC", str(value or ""))


def stable_id(*parts: Any, length: int = 32) -> str:
    raw = "::".join(nfc(part) for part in parts)
    return hashlib.sha256(raw.encode("utf-8")).hexdigest()[:length]


def clean_text(value: Any) -> str:
    text = nfc(value).replace("\x00", " ")
    text = re.sub(r"[ \t\r\f\v]+", " ", text)
    text = re.sub(r"\n{3,}", "\n\n", text)
    lines = [line.strip() for line in text.splitlines()]
    return "\n".join(line for line in lines if line).strip()


def text_quality(text: str, confidence_score: float | None = None) -> str:
    content = clean_text(text)
    length = len(content)
    if confidence_score is not None and confidence_score < 0.45:
        return "low"
    if length >= 500 and (confidence_score is None or confidence_score >= 0.70):
        return "high"
    if length >= 120:
        return "medium"
    return "low"


def split_text(text: str, max_chars: int, overlap_chars: int) -> list[str]:
    text = clean_text(text)
    if not text:
        return []
    paragraphs = re.split(r"\n\s*\n", text)
    chunks: list[str] = []
    current = ""
    for paragraph in paragraphs:
        paragraph = clean_text(paragraph)
        if not paragraph:
            continue
        if len(paragraph) > max_chars:
            if current:
                chunks.append(current)
                current = ""
            start = 0
            while start < len(paragraph):
                end = min(start + max_chars, len(paragraph))
                chunks.append(paragraph[start:end].strip())
                if end == len(paragraph):
                    break
                start = max(0, end - overlap_chars)
            continue
        candidate = clean_text(current + "\n\n" + paragraph) if current else paragraph
        if len(candidate) <= max_chars:
            current = candidate
            continue
        if current:
            chunks.append(current)
        prefix = current[-overlap_chars:] if overlap_chars and current else ""
        current = clean_text(prefix + "\n\n" + paragraph) if prefix else paragraph
    if current:
        chunks.append(current)
    return chunks


def selected_indices(total_count: int, args: argparse.Namespace) -> list[int]:
    if total_count <= 0:
        return []
    if args.sample_pages > 0:
        sample_count = min(args.sample_pages, total_count)
        if sample_count == total_count:
            return list(range(total_count))
        # Avoid front matter when possible. Most textbooks place covers, prefaces,
        # and table of contents in the first pages, so sampling starts after page 10
        # or around 20% of the document, whichever is smaller.
        start = min(max(10, int(total_count * 0.2)), max(0, total_count - sample_count))
        if sample_count == 1:
            return [min(total_count - 1, max(start, total_count // 2))]
        end = total_count - 1
        if end <= start:
            return list(range(sample_count))
        indices = []
        for i in range(sample_count):
            pos = round(start + (end - start) * (i / (sample_count - 1)))
            indices.append(min(total_count - 1, max(0, int(pos))))
        return sorted(set(indices))
    if args.max_pages:
        return list(range(min(args.max_pages, total_count)))
    return list(range(total_count))


def relative_to_project(path: Path) -> str:
    try:
        return str(Path("materials") / path.relative_to(EXTERNAL_MATERIALS_ROOT))
    except ValueError:
        pass
    try:
        return str(path.relative_to(PROJECT_ROOT))
    except ValueError:
        return str(path)


def material_folder(path: Path, input_dir: Path) -> str:
    try:
        rel = path.relative_to(input_dir)
    except ValueError:
        try:
            rel = path.relative_to(DEFAULT_INPUT.resolve())
        except ValueError:
            try:
                rel = path.relative_to(EXTERNAL_MATERIALS_ROOT / "04_subject_references")
            except ValueError:
                return ""
    return str(rel.parent)


def common_chunk_fields(
    *,
    chunk_id: str,
    document_id: str,
    chunk_type: str,
    source_file: str,
    source_path: str,
    page_or_slide: int,
    bbox: list[float] | None,
    content: str,
    structured_content: dict[str, Any] | None,
    source_image_path: str,
    extraction_method: str,
    extraction_quality: str,
    confidence_score: float,
    needs_review: bool,
    review_reason: str,
    vision_generated: bool = False,
    created_at: str,
) -> dict[str, Any]:
    record = {
        "chunk_id": chunk_id,
        "document_id": document_id,
        "chunk_type": chunk_type,
        "source_file": source_file,
        "source_path": source_path,
        "page_or_slide": page_or_slide,
        "bbox": bbox or [],
        "content": content,
        "structured_content": structured_content or {},
        "source_image_path": source_image_path,
        "extraction_method": extraction_method,
        "extraction_quality": extraction_quality,
        "confidence_score": round(float(confidence_score), 4),
        "needs_review": bool(needs_review),
        "review_reason": review_reason,
        "vision_generated": bool(vision_generated),
        "linked_scope": "",
        "linked_learning_objective": "",
        "scope_link_confidence": 0.0,
        "objective_link_confidence": 0.0,
        "linking_needs_review": True,
        "approved_for_generation": False,
        "pre_mapping_generation_candidate": chunk_type == "text" and extraction_quality == "high" and not needs_review,
        "created_at": created_at,
    }
    if chunk_type in {"table", "formula", "figure", "diagram"}:
        record["structured_content"].setdefault("embedded_text_candidates", [])
        record["embedded_text_candidates"] = record["structured_content"].get("embedded_text_candidates", [])
        record["needs_vision_model"] = chunk_type in {"figure", "diagram"}
        record["multimodal_interpretation"] = {
            "status": "pending",
            "semantic_description": "",
            "key_concepts": [],
            "related_rag_queries": [],
            "canonical_representation": {
                "formula_latex": "",
                "formula_plain_text": "",
                "table_json": {},
                "table_markdown": "",
                "diagram_spec": "",
                "chart_data": {},
            },
            "source_crosscheck": {
                "status": "pending",
                "supporting_chunk_ids": [],
                "supporting_source_pages": [],
                "confidence_score": 0.0,
                "notes": "",
            },
            "reconstruction_prompt": "",
            "generation_use_policy": (
                "Do not copy source visuals verbatim. Use the preserved crop, caption, "
                "nearby text, and RAG cross-check to create a new equivalent question "
                "asset only after review."
            ),
            "can_be_used_for_generation": False,
            "can_be_reconstructed": False,
            "requires_vision_llm": chunk_type in {"figure", "diagram"},
            "requires_human_review": True,
        }
    return record


def preprocess_for_ocr(image: Image.Image) -> Image.Image:
    if cv2 is None or np is None:
        return image
    array = np.array(image.convert("RGB"))
    gray = cv2.cvtColor(array, cv2.COLOR_RGB2GRAY)
    gray = cv2.medianBlur(gray, 3)
    thresholded = cv2.adaptiveThreshold(
        gray,
        255,
        cv2.ADAPTIVE_THRESH_GAUSSIAN_C,
        cv2.THRESH_BINARY,
        35,
        11,
    )
    return Image.fromarray(thresholded)


def padded_clip(page: fitz.Page, clip: fitz.Rect | None, padding: float = 0) -> fitz.Rect | None:
    if clip is None:
        return None
    return fitz.Rect(
        max(page.rect.x0, clip.x0 - padding),
        max(page.rect.y0, clip.y0 - padding),
        min(page.rect.x1, clip.x1 + padding),
        min(page.rect.y1, clip.y1 + padding),
    )


def render_clip(page: fitz.Page, clip: fitz.Rect | None, output_path: Path, dpi: int, padding: float = 4) -> str:
    output_path.parent.mkdir(parents=True, exist_ok=True)
    matrix = fitz.Matrix(dpi / 72, dpi / 72)
    pixmap = page.get_pixmap(matrix=matrix, clip=padded_clip(page, clip, padding), alpha=False)
    pixmap.save(str(output_path))
    return str(output_path)


def embedded_text_candidates_from_image(image_path: str, lang: str = "kor+eng", timeout: int = 6) -> list[str]:
    """OCR a crop as multimodal metadata, never as ordinary body text."""
    if not image_path or pytesseract is None:
        return []
    try:
        image = Image.open(image_path)
        text = pytesseract.image_to_string(image, lang=lang, config="--psm 6", timeout=timeout)
    except Exception:
        return []
    candidates: list[str] = []
    seen: set[str] = set()
    for line in clean_text(text).splitlines():
        line = clean_text(line)
        if not line or len(line) > 160:
            continue
        for piece in re.split(r"\s{2,}|[|]+", line):
            candidate = clean_text(piece)
            if not candidate or len(candidate) > 80 or candidate in seen:
                continue
            seen.add(candidate)
            candidates.append(candidate)
            if len(candidates) >= 30:
                return candidates
    return candidates


def refine_embedded_text_candidates(
    candidates: list[str],
    *,
    caption: str = "",
    nearby_text: str = "",
    max_length: int = 40,
) -> list[str]:
    reference_text = clean_text("\n".join([caption, nearby_text]))
    refined: list[str] = []
    seen: set[str] = set()
    for raw in candidates:
        candidate = clean_text(raw)
        if not candidate or len(candidate) > max_length:
            continue
        if re.search(r"(그림|표)\s*\d", candidate):
            continue
        if re.match(r"^\d+\s+.*학$", candidate):
            continue
        if candidate.endswith(("다.", "이다.", "한다.", "있다.", "된다.")):
            continue
        noisy_chars = sum(candidate.count(ch) for ch in ["\\", "`", "~", "ㄴ", "ㅁ"])
        if noisy_chars >= 5:
            continue
        if len(candidate) > 8 and candidate in reference_text:
            continue
        if re.fullmatch(r"[가-힣\s]{14,}", candidate):
            continue
        if candidate in seen:
            continue
        seen.add(candidate)
        refined.append(candidate)
    return refined


def ocr_pdf_page(page: fitz.Page, dpi: int, lang: str, timeout: int) -> tuple[str, float | None, str]:
    if pytesseract is None:
        return "", None, "pytesseract_not_installed"
    matrix = fitz.Matrix(dpi / 72, dpi / 72)
    pixmap = page.get_pixmap(matrix=matrix, alpha=False)
    image = Image.open(io.BytesIO(pixmap.tobytes("png")))
    image = preprocess_for_ocr(image)
    try:
        text = pytesseract.image_to_string(image, lang=lang, config="--psm 6", timeout=timeout)
    except RuntimeError as exc:
        return "", None, f"ocr_timeout_or_runtime_error:{exc}"
    values: list[float] = []
    try:
        data = pytesseract.image_to_data(
            image,
            lang=lang,
            config="--psm 6",
            output_type=pytesseract.Output.DICT,
            timeout=timeout,
        )
        for raw in data.get("conf", []):
            try:
                value = float(raw)
            except Exception:
                continue
            if value >= 0:
                values.append(value / 100)
    except Exception:
        pass
    confidence = sum(values) / len(values) if values else None
    return clean_text(text), confidence, ""


def table_to_markdown(table: list[list[Any]]) -> str:
    rows = []
    for row in table or []:
        cells = [clean_text(cell) for cell in row]
        if any(cells):
            rows.append(cells)
    if not rows:
        return ""
    width = max(len(row) for row in rows)
    rows = [row + [""] * (width - len(row)) for row in rows]
    header = rows[0]
    body = rows[1:]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * width) + " |"]
    for row in body:
        lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


def table_quality(table: list[list[Any]], markdown: str) -> tuple[str, float, str]:
    row_count = len(table or [])
    column_count = max((len(row or []) for row in table or []), default=0)
    non_empty_cells = sum(1 for row in table or [] for cell in row or [] if clean_text(cell))
    if row_count >= 2 and column_count >= 2 and non_empty_cells >= 4 and len(markdown) >= 80:
        return "high", 0.82, "structured_table_extracted"
    if row_count >= 1 and column_count >= 2 and non_empty_cells >= 2:
        return "medium", 0.62, "table_structure_partial"
    return "low", 0.35, "table_structure_uncertain"


def multimodal_seed(
    *,
    chunk_type: str,
    caption: str,
    nearby_text: str,
    table_markdown: str = "",
    table_json: list[list[Any]] | None = None,
    formula_text: str = "",
) -> dict[str, Any]:
    base_text = clean_text("\n".join([caption, nearby_text, table_markdown, formula_text]))
    terms = []
    for token in re.findall(r"[A-Za-z가-힣0-9]+(?:[·ㆍ/-][A-Za-z가-힣0-9]+)*", base_text):
        if len(token) >= 2 and token not in terms:
            terms.append(token)
        if len(terms) >= 12:
            break

    if chunk_type == "table":
        semantic_hint = "표의 행과 열이 비교하는 항목, 단위, 조건을 주변 문맥과 함께 해석해야 한다."
        reconstruction_hint = "원본 표를 복사하지 말고 같은 개념을 새로운 표나 조건 비교 문항으로 재구성한다."
    elif chunk_type == "formula":
        semantic_hint = "수식의 변수, 단위, 적용 조건을 주변 문맥과 원본 수식 이미지로 확인해야 한다."
        reconstruction_hint = "원본 수식 이미지를 확인한 뒤 같은 관계식을 새 문항에 맞게 LaTeX 또는 평문 수식으로 재작성한다."
    else:
        semantic_hint = "그림/도식의 축, 화살표, 구성 요소, 캡션을 주변 문맥과 함께 해석해야 한다."
        reconstruction_hint = "원본 그림을 복사하지 말고 같은 원리를 설명하는 새 도식 또는 그래프를 생성한다."

    return {
        "semantic_hint": semantic_hint,
        "keyword_candidates": terms,
        "related_rag_query_candidates": [
            clean_text(" ".join(terms[:5])),
            clean_text(caption),
        ],
        "canonical_representation_seed": {
            "formula_plain_text": formula_text,
            "table_markdown": table_markdown,
            "table_json": table_json or [],
        },
        "reconstruction_prompt_seed": reconstruction_hint,
    }


def extract_tables_for_page(
    *,
    pdf_path: Path,
    fitz_doc: fitz.Document,
    pdfplumber_page: Any,
    page_number: int,
    document_id: str,
    asset_dir: Path,
    created_at: str,
    render_dpi: int,
    save_crops: bool,
) -> list[dict[str, Any]]:
    chunks: list[dict[str, Any]] = []
    try:
        found_tables = pdfplumber_page.find_tables()
    except Exception:
        found_tables = []

    for table_index, table_obj in enumerate(found_tables, start=1):
        try:
            cells = table_obj.extract()
        except Exception:
            continue
        markdown = table_to_markdown(cells)
        if not markdown:
            continue
        bbox = [round(float(x), 2) for x in table_obj.bbox]
        quality, confidence, reason = table_quality(cells, markdown)
        visual_context = visual_context_for_page(fitz_doc[page_number - 1], bbox)
        source_image_path = ""
        if save_crops:
            asset_id = stable_id(pdf_path.name, page_number, "table", table_index, bbox)
            image_path = asset_dir / "tables" / f"{asset_id}.png"
            source_image_path = render_clip(fitz_doc[page_number - 1], fitz.Rect(*bbox), image_path, render_dpi)
        embedded_text_candidates = refine_embedded_text_candidates(
            embedded_text_candidates_from_image(source_image_path),
            caption=visual_context["caption"],
            nearby_text=visual_context["nearby_text"],
            max_length=60,
        )
        chunk_id = stable_id(document_id, page_number, "table", table_index, markdown)
        chunks.append(common_chunk_fields(
            chunk_id=chunk_id,
            document_id=document_id,
            chunk_type="table",
            source_file=pdf_path.name,
            source_path=relative_to_project(pdf_path),
            page_or_slide=page_number,
            bbox=bbox,
            content=markdown,
            structured_content={
                "table_json": cells,
                "table_markdown": markdown,
                "row_count": len(cells or []),
                "column_count": max((len(row or []) for row in cells or []), default=0),
                "table_index": table_index,
                "caption": visual_context["caption"],
                "context_before": visual_context["context_before"],
                "context_after": visual_context["context_after"],
                "nearby_text": visual_context["nearby_text"],
                "embedded_text_candidates": embedded_text_candidates,
                "multimodal_seed": multimodal_seed(
                    chunk_type="table",
                    caption=visual_context["caption"],
                    nearby_text=visual_context["nearby_text"],
                    table_markdown=markdown,
                    table_json=cells,
                ),
            },
            source_image_path=source_image_path,
            extraction_method="pdfplumber_find_tables",
            extraction_quality=quality,
            confidence_score=confidence,
            needs_review=quality != "high",
            review_reason=reason if quality != "high" else "",
            created_at=created_at,
        ))
    return chunks


def iter_text_lines(page: fitz.Page) -> list[dict[str, Any]]:
    lines: list[dict[str, Any]] = []
    data = page.get_text("dict", sort=True)
    for block in data.get("blocks", []):
        for line in block.get("lines", []):
            spans = line.get("spans", [])
            text = clean_text(" ".join(span.get("text", "") for span in spans))
            if not text:
                continue
            bbox = line.get("bbox", block.get("bbox", []))
            lines.append({"text": text, "bbox": [round(float(x), 2) for x in bbox]})
    return lines


def line_y0(line: dict[str, Any]) -> float:
    bbox = line.get("bbox") or [0, 0, 0, 0]
    return float(bbox[1]) if len(bbox) >= 2 else 0.0


def line_y1(line: dict[str, Any]) -> float:
    bbox = line.get("bbox") or [0, 0, 0, 0]
    return float(bbox[3]) if len(bbox) >= 4 else 0.0


def visual_context_for_page(page: fitz.Page, bbox: list[float] | None) -> dict[str, Any]:
    """Attach nearby prose/caption so multimodal chunks keep their meaning.

    Tables, equations, and diagrams often become ambiguous if only their crop
    image is stored. This context is used later by a Vision LLM or reviewer to
    describe what the visual item explains.
    """
    lines = iter_text_lines(page)
    if not lines:
        return {
            "caption": "",
            "context_before": "",
            "context_after": "",
            "nearby_text": "",
        }

    caption_pattern = re.compile(r"(?:그림|표|수식)\s*\d|핵심정리")
    caption_lines = [line for line in lines if caption_pattern.search(line["text"])]
    if bbox and len(bbox) >= 4:
        top = float(bbox[1])
        bottom = float(bbox[3])
        before = [line for line in lines if line_y1(line) <= top]
        after = [line for line in lines if line_y0(line) >= bottom]
        inside = [line for line in lines if line_y0(line) >= top and line_y1(line) <= bottom]
        nearby_captions = [
            line for line in caption_lines
            if abs(line_y0(line) - bottom) < 80 or abs(line_y1(line) - top) < 80
        ]
        caption = clean_text("\n".join(line["text"] for line in nearby_captions[:2]))
        context_before = clean_text("\n".join(line["text"] for line in before[-4:]))
        context_after = clean_text("\n".join(line["text"] for line in after[:6]))
        nearby_text = clean_text("\n".join(
            [context_before, caption, "\n".join(line["text"] for line in inside[:8]), context_after]
        ))
    else:
        caption = clean_text("\n".join(line["text"] for line in caption_lines[:4]))
        context_before = ""
        context_after = ""
        nearby_text = clean_text("\n".join(line["text"] for line in lines[:14]))

    return {
        "caption": caption,
        "context_before": context_before,
        "context_after": context_after,
        "nearby_text": nearby_text[:2400],
    }


def is_formula_candidate(text: str) -> bool:
    text = clean_text(text)
    if len(text) < 4 or len(text) > 220:
        return False
    if re.search(r"(전화|FAX|ISBN|등록번호|주소|http|www\.|@)", text, re.IGNORECASE):
        return False
    if not any(pattern.search(text) for pattern in FORMULA_PATTERNS):
        return False
    if re.fullmatch(r"[\w가-힣/%·().,\s-]+", text) and not re.search(r"[=＝∑√πΩμθλαβγ∆Δ±≤≥≠∞∝×÷*/+\ue000-\uf8ff]", text):
        return False
    if re.fullmatch(r"[A-Za-z가-힣/%·().,\s-]+", text):
        return False
    return True


def extract_formula_chunks(
    *,
    pdf_path: Path,
    page: fitz.Page,
    page_number: int,
    document_id: str,
    asset_dir: Path,
    created_at: str,
    render_dpi: int,
    save_crops: bool,
) -> list[dict[str, Any]]:
    chunks: list[dict[str, Any]] = []
    lines = iter_text_lines(page)
    for line_index, line in enumerate(lines):
        text = line["text"]
        if not is_formula_candidate(text):
            continue
        context = "\n".join(item["text"] for item in lines[max(0, line_index - 1): min(len(lines), line_index + 2)])
        bbox = line.get("bbox") or []
        visual_context = visual_context_for_page(page, bbox)
        source_image_path = ""
        if save_crops and bbox:
            asset_id = stable_id(pdf_path.name, page_number, "formula", line_index, text)
            image_path = asset_dir / "formulas" / f"{asset_id}.png"
            source_image_path = render_clip(page, fitz.Rect(*bbox), image_path, render_dpi)
        embedded_text_candidates = refine_embedded_text_candidates(
            embedded_text_candidates_from_image(source_image_path),
            caption=visual_context["caption"],
            nearby_text=visual_context["nearby_text"],
            max_length=60,
        )
        chunk_id = stable_id(document_id, page_number, "formula", line_index, text)
        chunks.append(common_chunk_fields(
            chunk_id=chunk_id,
            document_id=document_id,
            chunk_type="formula",
            source_file=pdf_path.name,
            source_path=relative_to_project(pdf_path),
            page_or_slide=page_number,
            bbox=bbox,
            content=f"{text}\n\n[주변 문맥]\n{context}",
            structured_content={
                "formula_text": text,
                "context": context,
                "line_index": line_index,
                "caption": visual_context["caption"],
                "context_before": visual_context["context_before"],
                "context_after": visual_context["context_after"],
                "nearby_text": visual_context["nearby_text"],
                "embedded_text_candidates": embedded_text_candidates,
                "multimodal_seed": multimodal_seed(
                    chunk_type="formula",
                    caption=visual_context["caption"],
                    nearby_text=visual_context["nearby_text"],
                    formula_text=text,
                ),
            },
            source_image_path=source_image_path,
            extraction_method="pymupdf_formula_heuristic",
            extraction_quality="medium",
            confidence_score=0.55,
            needs_review=True,
            review_reason="formula_candidate_requires_manual_review",
            created_at=created_at,
        ))
    return chunks


def image_area_ratio(page: fitz.Page, bbox: fitz.Rect) -> float:
    page_area = page.rect.width * page.rect.height
    return (bbox.width * bbox.height) / page_area if page_area else 0.0


def extract_visual_chunks(
    *,
    pdf_path: Path,
    page: fitz.Page,
    page_number: int,
    document_id: str,
    asset_dir: Path,
    created_at: str,
    render_dpi: int,
    min_area_ratio: float,
    max_area_ratio: float,
    drawing_threshold: int,
    save_crops: bool,
    max_visual_assets: int,
    current_visual_count: int,
) -> tuple[list[dict[str, Any]], int]:
    chunks: list[dict[str, Any]] = []
    drawings = len(page.get_drawings())
    image_infos = page.get_image_info(xrefs=True)
    text_lines = iter_text_lines(page)

    # Many textbook diagrams are vector drawings, not embedded image objects.
    # In those cases the most reliable anchor is the caption; crop the visual
    # area immediately above "그림 n-n" so the image, caption, and prose can be
    # interpreted together later.
    figure_caption_pattern = re.compile(r"^그림\s*\d+\s*[-–]\s*\d+\.")
    figure_caption_lines = [
        line for line in text_lines
        if figure_caption_pattern.search(line["text"])
    ]
    for caption_index, caption_line in enumerate(figure_caption_lines, start=1):
        if current_visual_count >= max_visual_assets:
            break
        caption_bbox = caption_line.get("bbox") or []
        if len(caption_bbox) < 4:
            continue
        caption_top = float(caption_bbox[1])
        previous_caption_bottoms = [
            line_y1(line)
            for line in figure_caption_lines
            if line_y1(line) < caption_top
        ]
        previous_boundary = previous_caption_bottoms[-1] + 16 if previous_caption_bottoms else 40
        top_limit = max(previous_boundary, caption_top - 340)
        bottom_limit = max(top_limit + 20, line_y1(caption_line) + 8)
        bbox = fitz.Rect(40, top_limit, page.rect.width - 40, bottom_limit)
        if bbox.height < 35 or bbox.width < 80:
            continue
        visual_context = visual_context_for_page(page, [bbox.x0, bbox.y0, bbox.x1, bbox.y1])
        caption = clean_text(caption_line.get("text", ""))
        if caption and caption not in visual_context["caption"]:
            visual_context["caption"] = clean_text("\n".join([visual_context["caption"], caption]))
        source_image_path = ""
        if save_crops:
            asset_id = stable_id(pdf_path.name, page_number, "caption_figure", caption_index, bbox)
            image_path = asset_dir / "figures" / f"{asset_id}.png"
            source_image_path = render_clip(page, bbox, image_path, render_dpi, padding=8)
        embedded_text_candidates = refine_embedded_text_candidates(
            embedded_text_candidates_from_image(source_image_path),
            caption=visual_context["caption"],
            nearby_text=visual_context["nearby_text"],
        )
        chunk_id = stable_id(document_id, page_number, "caption_figure", caption_index, caption)
        chunks.append(common_chunk_fields(
            chunk_id=chunk_id,
            document_id=document_id,
            chunk_type="figure",
            source_file=pdf_path.name,
            source_path=relative_to_project(pdf_path),
            page_or_slide=page_number,
            bbox=[round(float(x), 2) for x in [bbox.x0, bbox.y0, bbox.x1, bbox.y1]],
            content="[Vision LLM 설명 대기] 그림 캡션 기준으로 crop한 시각 자료 후보입니다.",
            structured_content={
                "caption_index": caption_index,
                "caption": visual_context["caption"],
                "context_before": visual_context["context_before"],
                "context_after": visual_context["context_after"],
                "nearby_text": visual_context["nearby_text"],
                "embedded_text_candidates": embedded_text_candidates,
                "description_status": "needs_vision_llm",
                "multimodal_seed": multimodal_seed(
                    chunk_type="figure",
                    caption=visual_context["caption"],
                    nearby_text=visual_context["nearby_text"],
                ),
            },
            source_image_path=source_image_path,
            extraction_method="pymupdf_caption_figure_crop",
            extraction_quality="medium",
            confidence_score=0.68,
            needs_review=True,
            review_reason="caption_figure_requires_vision_or_human_review",
            created_at=created_at,
        ))
        current_visual_count += 1

    for image_index, info in enumerate(image_infos, start=1):
        if current_visual_count >= max_visual_assets:
            break
        bbox = fitz.Rect(info.get("bbox", (0, 0, 0, 0)))
        ratio = image_area_ratio(page, bbox)
        if ratio < min_area_ratio or ratio > max_area_ratio:
            continue
        visual_context = visual_context_for_page(page, [bbox.x0, bbox.y0, bbox.x1, bbox.y1])
        source_image_path = ""
        if save_crops:
            asset_id = stable_id(pdf_path.name, page_number, "image", image_index, bbox)
            image_path = asset_dir / "figures" / f"{asset_id}.png"
            source_image_path = render_clip(page, bbox, image_path, render_dpi)
        embedded_text_candidates = refine_embedded_text_candidates(
            embedded_text_candidates_from_image(source_image_path),
            caption=visual_context["caption"],
            nearby_text=visual_context["nearby_text"],
        )
        chunk_id = stable_id(document_id, page_number, "figure", image_index, bbox)
        chunks.append(common_chunk_fields(
            chunk_id=chunk_id,
            document_id=document_id,
            chunk_type="figure",
            source_file=pdf_path.name,
            source_path=relative_to_project(pdf_path),
            page_or_slide=page_number,
            bbox=[round(float(x), 2) for x in [bbox.x0, bbox.y0, bbox.x1, bbox.y1]],
            content="[Vision LLM 설명 대기] 원본 이미지 crop을 보존한 시각 자료 후보입니다.",
            structured_content={
                "image_index": image_index,
                "page_area_ratio": round(ratio, 4),
                "drawing_count": drawings,
                "description_status": "needs_vision_llm",
                "caption": visual_context["caption"],
                "context_before": visual_context["context_before"],
                "context_after": visual_context["context_after"],
                "nearby_text": visual_context["nearby_text"],
                "embedded_text_candidates": embedded_text_candidates,
                "multimodal_seed": multimodal_seed(
                    chunk_type="figure",
                    caption=visual_context["caption"],
                    nearby_text=visual_context["nearby_text"],
                ),
            },
            source_image_path=source_image_path,
            extraction_method="pymupdf_image_clip",
            extraction_quality="medium",
            confidence_score=0.55,
            needs_review=True,
            review_reason="figure_requires_vision_or_human_review",
            created_at=created_at,
        ))
        current_visual_count += 1

    if drawings >= drawing_threshold and current_visual_count < max_visual_assets:
        visual_context = visual_context_for_page(page, None)
        source_image_path = ""
        if save_crops:
            asset_id = stable_id(pdf_path.name, page_number, "diagram_page", drawings)
            image_path = asset_dir / "figures" / f"{asset_id}.png"
            source_image_path = render_clip(page, None, image_path, render_dpi)
        embedded_text_candidates = []
        chunk_id = stable_id(document_id, page_number, "diagram_page", drawings)
        chunks.append(common_chunk_fields(
            chunk_id=chunk_id,
            document_id=document_id,
            chunk_type="diagram",
            source_file=pdf_path.name,
            source_path=relative_to_project(pdf_path),
            page_or_slide=page_number,
            bbox=[],
            content="[Vision LLM 설명 대기] 벡터 도식/장비 구조도 가능성이 있는 페이지입니다.",
            structured_content={
                "drawing_count": drawings,
                "description_status": "needs_vision_llm",
                "caption": visual_context["caption"],
                "context_before": visual_context["context_before"],
                "context_after": visual_context["context_after"],
                "nearby_text": visual_context["nearby_text"],
                "embedded_text_candidates": embedded_text_candidates,
                "multimodal_seed": multimodal_seed(
                    chunk_type="diagram",
                    caption=visual_context["caption"],
                    nearby_text=visual_context["nearby_text"],
                ),
            },
            source_image_path=source_image_path,
            extraction_method="pymupdf_page_render_drawing_candidate",
            extraction_quality="medium",
            confidence_score=0.50,
            needs_review=True,
            review_reason="diagram_requires_vision_or_human_review",
            created_at=created_at,
        ))
        current_visual_count += 1

    return chunks, current_visual_count


def extract_pdf(path: Path, args: argparse.Namespace, created_at: str) -> tuple[dict[str, Any], list[dict[str, Any]], list[dict[str, Any]]]:
    document_id = stable_id(relative_to_project(path), path.stat().st_size, path.stat().st_mtime_ns)
    chunks: list[dict[str, Any]] = []
    page_records: list[dict[str, Any]] = []
    visual_count = 0
    fitz_doc = fitz.open(str(path))
    page_indices = selected_indices(len(fitz_doc), args)

    plumber_pdf = None
    try:
        plumber_pdf = pdfplumber.open(str(path))
    except Exception:
        plumber_pdf = None

    for page_index in page_indices:
        page_number = page_index + 1
        page = fitz_doc[page_index]
        text = clean_text(page.get_text("text", sort=True))
        method = "pymupdf_text"
        ocr_confidence = None
        ocr_note = ""
        if args.ocr and len(text) < args.ocr_min_chars:
            ocr_text, ocr_confidence, ocr_note = ocr_pdf_page(page, args.ocr_dpi, args.ocr_lang, args.ocr_timeout)
            if len(ocr_text) > len(text):
                text = ocr_text
                method = "tesseract_ocr"
            elif ocr_text:
                text = clean_text(text + "\n\n[OCR 보조 추출]\n" + ocr_text)
                method = "pymupdf_text+tesseract_ocr"

        confidence = ocr_confidence if ocr_confidence is not None else (0.92 if text else 0.0)
        quality = text_quality(text, confidence)
        page_review_reason = "" if quality == "high" else f"text_quality_{quality}"
        text_chunk_count = 0
        for chunk_index, content in enumerate(split_text(text, args.chunk_chars, args.overlap_chars)):
            chunk_id = stable_id(document_id, page_number, "text", chunk_index, content)
            chunks.append(common_chunk_fields(
                chunk_id=chunk_id,
                document_id=document_id,
                chunk_type="text",
                source_file=path.name,
                source_path=relative_to_project(path),
                page_or_slide=page_number,
                bbox=[],
                content=content,
                structured_content={
                    "chunk_index": chunk_index,
                    "token_estimate": max(1, len(content) // 3),
                },
                source_image_path="",
                extraction_method=method,
                extraction_quality=quality,
                confidence_score=confidence,
                needs_review=quality != "high",
                review_reason=page_review_reason,
                created_at=created_at,
            ))
            text_chunk_count += 1

        table_chunks: list[dict[str, Any]] = []
        if plumber_pdf and args.extract_tables and page_index < len(plumber_pdf.pages):
            table_chunks = extract_tables_for_page(
                pdf_path=path,
                fitz_doc=fitz_doc,
                pdfplumber_page=plumber_pdf.pages[page_index],
                page_number=page_number,
                document_id=document_id,
                asset_dir=args.asset_output,
                created_at=created_at,
                render_dpi=args.asset_dpi,
                save_crops=args.save_crops,
            )
            chunks.extend(table_chunks)

        formula_chunks: list[dict[str, Any]] = []
        if args.extract_formulas:
            formula_chunks = extract_formula_chunks(
                pdf_path=path,
                page=page,
                page_number=page_number,
                document_id=document_id,
                asset_dir=args.asset_output,
                created_at=created_at,
                render_dpi=args.asset_dpi,
                save_crops=args.save_crops,
            )
            chunks.extend(formula_chunks)

        visual_chunks: list[dict[str, Any]] = []
        if args.extract_visuals:
            visual_chunks, visual_count = extract_visual_chunks(
                pdf_path=path,
                page=page,
                page_number=page_number,
                document_id=document_id,
                asset_dir=args.asset_output,
                created_at=created_at,
                render_dpi=args.asset_dpi,
                min_area_ratio=args.visual_image_min_area_ratio,
                max_area_ratio=args.visual_image_max_area_ratio,
                drawing_threshold=args.visual_drawing_threshold,
                save_crops=args.save_crops,
                max_visual_assets=args.max_visual_assets_per_document,
                current_visual_count=visual_count,
            )
            chunks.extend(visual_chunks)

        page_records.append({
            "page_or_slide": page_number,
            "text_chars": len(text),
            "text_quality": quality,
            "text_chunk_count": text_chunk_count,
            "table_chunk_count": len(table_chunks),
            "formula_chunk_count": len(formula_chunks),
                "visual_chunk_count": len(visual_chunks),
                "extraction_method": method,
                "ocr_confidence": ocr_confidence,
                "ocr_note": ocr_note,
            })

    if plumber_pdf:
        plumber_pdf.close()
    fitz_doc.close()

    document = {
        "document_id": document_id,
        "source_file": path.name,
        "source_path": relative_to_project(path),
        "source_type": "pdf",
        "material_folder": material_folder(path, args.input),
        "file_size": path.stat().st_size,
        "page_count": len(page_records),
        "created_at": created_at,
    }
    return document, chunks, page_records


def extract_text_like(path: Path, args: argparse.Namespace, created_at: str) -> tuple[dict[str, Any], list[dict[str, Any]], list[dict[str, Any]]]:
    suffix = path.suffix.lower()
    parts: list[str] = []
    source_type = suffix.lstrip(".")
    method = "plain_text"
    page_or_slide = 1
    if suffix in {".txt", ".md"}:
        parts.append(path.read_text(encoding="utf-8", errors="ignore"))
    elif suffix == ".docx":
        if Document is None:
            raise RuntimeError("python-docx is not installed")
        method = "python-docx"
        doc = Document(str(path))
        parts.extend(paragraph.text for paragraph in doc.paragraphs if clean_text(paragraph.text))
        for table in doc.tables:
            for row in table.rows:
                parts.append(" | ".join(clean_text(cell.text) for cell in row.cells))
    elif suffix == ".pptx":
        if Presentation is None:
            raise RuntimeError("python-pptx is not installed")
        method = "python-pptx"
        prs = Presentation(str(path))
        chunks: list[dict[str, Any]] = []
        page_records: list[dict[str, Any]] = []
        document_id = stable_id(relative_to_project(path), path.stat().st_size, path.stat().st_mtime_ns)
        slide_indices = selected_indices(len(prs.slides), args)
        for raw_index in slide_indices:
            slide = prs.slides[raw_index]
            slide_index = raw_index + 1
            slide_parts: list[str] = []
            table_count = 0
            image_count = 0
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_parts.append(shape.text)
                if getattr(shape, "has_table", False):
                    table_count += 1
                    for row in shape.table.rows:
                        slide_parts.append(" | ".join(clean_text(cell.text) for cell in row.cells))
                if getattr(shape, "shape_type", None) and "PICTURE" in str(shape.shape_type):
                    image_count += 1
            text = clean_text("\n".join(slide_parts))
            quality = text_quality(text)
            for chunk_index, content in enumerate(split_text(text, args.chunk_chars, args.overlap_chars)):
                chunk_id = stable_id(document_id, slide_index, "text", chunk_index, content)
                chunks.append(common_chunk_fields(
                    chunk_id=chunk_id,
                    document_id=document_id,
                    chunk_type="text",
                    source_file=path.name,
                    source_path=relative_to_project(path),
                    page_or_slide=slide_index,
                    bbox=[],
                    content=content,
                    structured_content={"chunk_index": chunk_index, "token_estimate": max(1, len(content) // 3)},
                    source_image_path="",
                    extraction_method=method,
                    extraction_quality=quality,
                    confidence_score=0.90 if quality == "high" else 0.60,
                    needs_review=quality != "high",
                    review_reason="" if quality == "high" else f"text_quality_{quality}",
                    created_at=created_at,
                ))
            page_records.append({
                "page_or_slide": slide_index,
                "text_chars": len(text),
                "text_quality": quality,
                "table_count": table_count,
                "image_count": image_count,
            })
        document = {
            "document_id": document_id,
            "source_file": path.name,
            "source_path": relative_to_project(path),
            "source_type": source_type,
            "material_folder": material_folder(path, args.input),
            "file_size": path.stat().st_size,
            "page_count": len(page_records),
            "created_at": created_at,
        }
        return document, chunks, page_records
    else:
        raise RuntimeError(f"unsupported file extension: {suffix}")

    text = clean_text("\n".join(parts))
    quality = text_quality(text)
    document_id = stable_id(relative_to_project(path), path.stat().st_size, path.stat().st_mtime_ns)
    chunks = []
    for chunk_index, content in enumerate(split_text(text, args.chunk_chars, args.overlap_chars)):
        chunk_id = stable_id(document_id, page_or_slide, "text", chunk_index, content)
        chunks.append(common_chunk_fields(
            chunk_id=chunk_id,
            document_id=document_id,
            chunk_type="text",
            source_file=path.name,
            source_path=relative_to_project(path),
            page_or_slide=page_or_slide,
            bbox=[],
            content=content,
            structured_content={"chunk_index": chunk_index, "token_estimate": max(1, len(content) // 3)},
            source_image_path="",
            extraction_method=method,
            extraction_quality=quality,
            confidence_score=0.90 if quality == "high" else 0.60,
            needs_review=quality != "high",
            review_reason="" if quality == "high" else f"text_quality_{quality}",
            created_at=created_at,
        ))
    document = {
        "document_id": document_id,
        "source_file": path.name,
        "source_path": relative_to_project(path),
        "source_type": source_type,
        "material_folder": material_folder(path, args.input),
        "file_size": path.stat().st_size,
        "page_count": 1,
        "created_at": created_at,
    }
    page_records = [{"page_or_slide": 1, "text_chars": len(text), "text_quality": quality}]
    return document, chunks, page_records


def discover_files(input_dir: Path) -> list[Path]:
    files = []
    for path in sorted(input_dir.rglob("*")):
        if not path.is_file():
            continue
        if path.name.startswith(".") or path.name.upper() == "README.MD":
            continue
        if path.suffix.lower() in SUPPORTED_EXTENSIONS:
            files.append(path)
    return files


def write_json(path: Path, data: Any) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    path.write_text(json.dumps(data, ensure_ascii=False, indent=2) + "\n", encoding="utf-8")


def write_jsonl(path: Path, rows: list[dict[str, Any]]) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)
    with path.open("w", encoding="utf-8") as f:
        for row in rows:
            f.write(json.dumps(row, ensure_ascii=False) + "\n")


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(description=__doc__)
    parser.add_argument("--input", type=Path, default=DEFAULT_INPUT)
    parser.add_argument("--output", type=Path, default=DEFAULT_OUTPUT)
    parser.add_argument("--asset-output", type=Path, default=DEFAULT_ASSET_OUTPUT)
    parser.add_argument("--chunk-chars", type=int, default=1400)
    parser.add_argument("--overlap-chars", type=int, default=180)
    parser.add_argument("--ocr-min-chars", type=int, default=120)
    parser.add_argument("--ocr-dpi", type=int, default=300)
    parser.add_argument("--ocr-lang", default="kor+eng")
    parser.add_argument("--ocr-timeout", type=int, default=30)
    parser.add_argument("--asset-dpi", type=int, default=180)
    parser.add_argument("--visual-image-min-area-ratio", type=float, default=0.15)
    parser.add_argument("--visual-image-max-area-ratio", type=float, default=0.9)
    parser.add_argument("--visual-drawing-threshold", type=int, default=500)
    parser.add_argument("--max-visual-assets-per-document", type=int, default=80)
    parser.add_argument("--max-pages", type=int, default=0, help="Limit pages/slides per file for smoke tests.")
    parser.add_argument("--sample-pages", type=int, default=0, help="Sample N non-front-matter pages/slides per file.")
    parser.add_argument("--no-ocr", dest="ocr", action="store_false")
    parser.add_argument("--no-tables", dest="extract_tables", action="store_false")
    parser.add_argument("--no-formulas", dest="extract_formulas", action="store_false")
    parser.add_argument("--no-visuals", dest="extract_visuals", action="store_false")
    parser.add_argument("--no-crops", dest="save_crops", action="store_false")
    parser.set_defaults(ocr=True, extract_tables=True, extract_formulas=True, extract_visuals=True, save_crops=True)
    return parser.parse_args()


def main() -> None:
    args = parse_args()
    args.input = args.input.resolve()
    args.output = args.output.resolve()
    args.asset_output = args.asset_output.resolve()
    created_at = now_iso()

    files = discover_files(args.input)
    documents: list[dict[str, Any]] = []
    chunks: list[dict[str, Any]] = []
    pages: list[dict[str, Any]] = []
    errors: list[dict[str, str]] = []

    for path in files:
        try:
            if path.suffix.lower() == ".pdf":
                document, file_chunks, page_records = extract_pdf(path, args, created_at)
            else:
                document, file_chunks, page_records = extract_text_like(path, args, created_at)
            documents.append(document)
            chunks.extend(file_chunks)
            for record in page_records:
                row = dict(record)
                row["document_id"] = document["document_id"]
                row["source_file"] = document["source_file"]
                row["source_path"] = document["source_path"]
                row["material_folder"] = document["material_folder"]
                pages.append(row)
        except Exception as exc:
            errors.append({
                "source_file": path.name,
                "source_path": relative_to_project(path),
                "error": f"{type(exc).__name__}: {exc}",
            })

    chunk_type_counts = Counter(row.get("chunk_type") for row in chunks)
    quality_counts = Counter(row.get("extraction_quality") for row in chunks)
    review_counts = Counter("needs_review" if row.get("needs_review") else "not_required" for row in chunks)
    report = {
        "version": 1,
        "created_at": created_at,
        "input_dir": str(args.input),
        "output_dir": str(args.output),
        "asset_output_dir": str(args.asset_output),
        "file_count": len(files),
        "document_count": len(documents),
        "page_count": len(pages),
        "chunk_count": len(chunks),
        "chunk_type_counts": dict(chunk_type_counts),
        "quality_counts": dict(quality_counts),
        "review_counts": dict(review_counts),
        "pre_mapping_generation_candidate_count": sum(1 for row in chunks if row.get("pre_mapping_generation_candidate")),
        "errors": errors,
        "notes": [
            "원본 파일은 수정하지 않는다.",
            "텍스트 high chunk만 pre_mapping_generation_candidate=true로 표시한다.",
            "표, 수식, 그림, 도식은 기본적으로 needs_review=true이며 자동 승인하지 않는다.",
            "Vision LLM 설명은 아직 생성하지 않고 source_image_path와 description_status 후보만 저장한다.",
            "이미지 crop은 /opt/app/extracted_assets 아래에 저장해 /mnt/data 용량을 보호한다.",
        ],
    }

    write_json(args.output / "documents.json", documents)
    write_jsonl(args.output / "chunks_all.jsonl", chunks)
    write_jsonl(args.output / "chunks_text_high.jsonl", [
        row for row in chunks
        if row.get("chunk_type") == "text" and row.get("extraction_quality") == "high" and not row.get("needs_review")
    ])
    write_jsonl(args.output / "review_queue.jsonl", [row for row in chunks if row.get("needs_review")])
    write_json(args.output / "pages.json", pages)
    write_json(args.output / "advanced_extraction_report.json", report)
    print(json.dumps(report, ensure_ascii=False, indent=2))


if __name__ == "__main__":
    main()
